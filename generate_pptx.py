#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Presentation PFE - ChatBot UNA v3.0
Design moderne, animations auto, liens cliquables
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree

# ======================================================
#  COULEURS
# ======================================================
G_DARK   = RGBColor(0x0d, 0x3d, 0x23)
G_MAIN   = RGBColor(0x1a, 0x5c, 0x35)
G_MED    = RGBColor(0x2d, 0x7a, 0x4f)
G_LIGHT  = RGBColor(0x4a, 0xde, 0x80)
G_PALE   = RGBColor(0xec, 0xfb, 0xf3)
G_TEXT   = RGBColor(0xbb, 0xf7, 0xd0)
WHITE    = RGBColor(0xff, 0xff, 0xff)
DARK     = RGBColor(0x0f, 0x17, 0x2a)
SLATE    = RGBColor(0x1e, 0x29, 0x3b)
GRAY     = RGBColor(0x64, 0x74, 0x8b)
LGRAY    = RGBColor(0xe2, 0xe8, 0xf0)
GOLD     = RGBColor(0xf5, 0x9e, 0x0b)
LGOLD    = RGBColor(0xff, 0xf7, 0xe0)
BLUE     = RGBColor(0x03, 0x69, 0xa1)
LBLUE    = RGBColor(0xdb, 0xea, 0xfe)
RED      = RGBColor(0xdc, 0x26, 0x26)
LRED     = RGBColor(0xfe, 0xe2, 0xe2)
ORANGE   = RGBColor(0xea, 0x58, 0x0c)
LORANGE  = RGBColor(0xff, 0xed, 0xd5)
PURPLE   = RGBColor(0x7e, 0x22, 0xce)
LPURPLE  = RGBColor(0xed, 0xe9, 0xfe)
BGLIGHT  = RGBColor(0xf8, 0xfa, 0xfc)

URL_CHATBOT = "https://chatbot-univn-production.up.railway.app/"
URL_ADMIN   = "https://chatbot-univn-production.up.railway.app/dashboard"

W = Inches(13.33)
H = Inches(7.5)
TOTAL = 18

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ======================================================
#  ANIMATIONS & TRANSITIONS
# ======================================================

def add_transition(sl, t="fade"):
    NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    try:
        xml = '<p:transition xmlns:p="{}" spd="med"><p:{}/></p:transition>'.format(NS, t)
        sl.element.append(etree.fromstring(xml))
    except Exception:
        pass


def add_auto_anim(slide, shape_ids, start_ms=150, step_ms=220, dur_ms=500):
    """Fade-in automatique (sans clic) - chaque forme apparait en fondu."""
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    if not shape_ids:
        return
    parts = ""
    for i, sp_id in enumerate(shape_ids):
        delay = start_ms + i * step_ms
        g = 100 + i * 10
        parts += """
        <p:par>
          <p:cTn id="{g}" presetID="10" presetClass="entr" presetSubtype="0"
                 fill="hold" grpId="{i}" nodeType="withPrevious">
            <p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>
            <p:childTnLst>
              <p:set>
                <p:cBhvr>
                  <p:cTn id="{g1}" dur="1" fill="hold"/>
                  <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                  <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                </p:cBhvr>
                <p:to><p:strVal val="visible"/></p:to>
              </p:set>
              <p:animEffect transition="in" filter="fade">
                <p:cBhvr>
                  <p:cTn id="{g2}" dur="{dur}"/>
                  <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                </p:cBhvr>
              </p:animEffect>
            </p:childTnLst>
          </p:cTn>
        </p:par>""".format(g=g, i=i, delay=delay, g1=g+1, g2=g+2, spid=sp_id, dur=dur_ms)

    timing = """<p:timing xmlns:p="{NP}" xmlns:a="{NA}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>
          <p:par>
            <p:cTn id="2" fill="hold">
              <p:stCondLst><p:cond delay="0" evt="begin"/></p:stCondLst>
              <p:childTnLst>{parts}</p:childTnLst>
            </p:cTn>
          </p:par>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>""".format(NP=NS_P, NA=NS_A, parts=parts)
    try:
        slide.element.append(etree.fromstring(timing))
    except Exception:
        pass


def get_content_ids(slide, skip=0):
    """Retourne shape_ids du contenu (en ignorant les skip premiers = fond/header/footer)."""
    all_ids = [s.shape_id for s in slide.shapes]
    return all_ids[skip:]


# ======================================================
#  HELPERS DE DESSIN
# ======================================================

def new_slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill=None, line=None, lw=Pt(1)):
    shp = s.shapes.add_shape(1, x, y, w, h)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width = lw
    else:
        shp.line.fill.background()
    return shp


def rnd(s, x, y, w, h, fill=None, line=None, lw=Pt(1)):
    shp = s.shapes.add_shape(5, x, y, w, h)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width = lw
    else:
        shp.line.fill.background()
    try:
        shp.adjustments[0] = 0.08
    except Exception:
        pass
    return shp


def txt(s, text, x, y, w, h, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def txt_link(s, text, url, x, y, w, h, size=12, color=BLUE, bold=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    try:
        r.hyperlink.address = url
    except Exception:
        pass
    return tb


def btn(s, label, url, x, y, w, h, bg=G_MAIN, fg=WHITE, size=13):
    shp = s.shapes.add_shape(5, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    shp.line.fill.background()
    try:
        shp.adjustments[0] = 0.18
    except Exception:
        pass
    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg
    r.font.name = "Calibri"
    try:
        r.hyperlink.address = url
    except Exception:
        pass
    return shp


def badge(s, x, y, w, h, text, bg, fg=WHITE, size=11, bold=True):
    shp = s.shapes.add_shape(5, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    shp.line.fill.background()
    try:
        shp.adjustments[0] = 0.3
    except Exception:
        pass
    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = fg
    r.font.name = "Calibri"
    return shp


def header(s, title, sub=None):
    rect(s, 0, 0, W, Inches(1.1), fill=G_MAIN)
    rect(s, 0, Inches(1.08), W, Inches(0.04), fill=G_LIGHT)
    txt(s, title, Inches(0.4), Inches(0.12), Inches(10), Inches(0.62),
        size=26, bold=True, color=WHITE)
    if sub:
        txt(s, sub, Inches(0.4), Inches(0.7), Inches(11), Inches(0.36),
            size=12, color=G_TEXT, italic=True)
    rect(s, W - Inches(0.5), Inches(0.35), Inches(0.08), Inches(0.08), fill=G_LIGHT)


def footer(s, num):
    rect(s, 0, H - Inches(0.38), W, Inches(0.38), fill=DARK)
    txt(s, "ChatBot UNA  |  Universite de Nouakchott  |  PFE 2024-2025",
        Inches(0.3), H - Inches(0.34), Inches(10), Inches(0.3),
        size=9, color=G_TEXT)
    txt(s, "{} / {}".format(num, TOTAL),
        W - Inches(1.2), H - Inches(0.34), Inches(0.9), Inches(0.3),
        size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, fill=WHITE, border=LGRAY):
    rect(s, x + Inches(0.055), y + Inches(0.055), w, h, fill=LGRAY)
    rnd(s, x, y, w, h, fill=fill, line=border, lw=Pt(1))


def sep(s, y, color=LGRAY):
    rect(s, Inches(0.4), y, W - Inches(0.8), Pt(2), fill=color)


def make_table(s, x, y, w, rows_data, col_widths,
               hdr_bg=G_MAIN, hdr_fg=WHITE,
               alt=G_PALE, normal=WHITE):
    rows = len(rows_data)
    cols = len(rows_data[0])
    tbl = s.shapes.add_table(rows, cols, x, y, w,
                              Inches(0.4) * rows).table
    tbl.first_row = True
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw
    for ri, row_data in enumerate(rows_data):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Calibri"
            if ri == 0:
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = hdr_fg
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_bg
            else:
                run.font.size = Pt(10)
                run.font.color.rgb = SLATE
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt if ri % 2 == 0 else normal
    return tbl


# ======================================================
#  SLIDE 1 — COUVERTURE
# ======================================================
s1 = new_slide()
rect(s1, 0, 0, W, H, fill=G_DARK)
for i in range(12):
    rect(s1, Inches(8.5 + i * 0.42), 0, Inches(0.28), H,
         fill=RGBColor(0x1a, 0x5c, 0x35))
rect(s1, 0, 0, Inches(8.4), Inches(4.4), fill=G_MAIN)
rect(s1, 0, Inches(4.36), Inches(8.4), Inches(0.06), fill=G_LIGHT)
# Contenu anime (skip 15 = 1 bg + 12 bandes + hero + ligne)
b_una  = badge(s1, Inches(0.45), Inches(0.35), Inches(1.1), Inches(0.48), "UNA", G_LIGHT, G_DARK, size=15)
t_main = txt(s1, "ChatBot Universitaire Intelligent",
             Inches(0.45), Inches(1.0), Inches(7.8), Inches(1.0),
             size=38, bold=True, color=WHITE)
t_sub  = txt(s1, "Universite de Nouakchott  -  Mauritanie",
             Inches(0.45), Inches(2.05), Inches(7.8), Inches(0.55),
             size=20, color=G_TEXT, italic=True)
t_tags = txt(s1, "Assistant bilingue (Arabe / Francais)  |  Rasa  |  Flask  |  React  |  Docker",
             Inches(0.45), Inches(2.7), Inches(7.8), Inches(0.4),
             size=13, color=G_LIGHT)
rect(s1, 0, Inches(4.42), W, H - Inches(4.42), fill=DARK)
info = [("Etudiants", "Med Mokhtar & Salim"),
        ("Encadrant", "Dr. [Nom]"),
        ("Annee",     "2024 - 2025"),
        ("Mention",   "Licence Informatique")]
cards_s1 = []
for i, (lbl, val) in enumerate(info):
    x = Inches(0.45) + i * Inches(2.05)
    c = rnd(s1, x, Inches(3.2), Inches(1.9), Inches(1.0),
            fill=RGBColor(0x12, 0x4a, 0x2c))
    cards_s1.append(c.shape_id)
    txt(s1, lbl.upper(), x + Inches(0.12), Inches(3.27), Inches(1.7), Inches(0.28),
        size=8, color=G_LIGHT, bold=True)
    txt(s1, val, x + Inches(0.12), Inches(3.55), Inches(1.7), Inches(0.38),
        size=12, color=WHITE, bold=True)
txt(s1, "Acces direct au projet :", Inches(0.5), Inches(4.65), Inches(5), Inches(0.4),
    size=14, bold=True, color=WHITE)
b1 = btn(s1, "  ChatBot en ligne  ", URL_CHATBOT,
         Inches(0.5), Inches(5.2), Inches(4.0), Inches(0.7),
         bg=G_LIGHT, fg=G_DARK, size=14)
b2 = btn(s1, "  Dashboard Admin  ", URL_ADMIN,
         Inches(4.9), Inches(5.2), Inches(4.0), Inches(0.7),
         bg=G_MED, fg=WHITE, size=14)
txt_link(s1, URL_CHATBOT, URL_CHATBOT,
         Inches(0.5), Inches(6.1), Inches(8), Inches(0.3),
         size=10, color=G_TEXT)
add_auto_anim(s1, [b_una.shape_id, t_main.shape_id, t_sub.shape_id,
                   t_tags.shape_id, b1.shape_id, b2.shape_id])
add_transition(s1, "fade")


# ======================================================
#  SLIDE 2 — SOMMAIRE
# ======================================================
s2 = new_slide()
rect(s2, 0, 0, W, H, fill=BGLIGHT)
header(s2, "Sommaire")
footer(s2, 2)

items = [
    ("01", "Contexte & Problematique",   "Mauritanie, UNA, besoin numerique"),
    ("02", "Etat de l'Art",              "Dialogflow vs Watson vs Rasa"),
    ("03", "Architecture 3-tiers",       "React  Flask  Rasa  SQLite  Docker"),
    ("04", "Pipeline NLU",               "DIETClassifier, politiques, 300+ intents"),
    ("05", "Diagramme de Sequence",      "Flux complet d'un message"),
    ("06", "Backend Flask",              "10 endpoints REST, WAL, fallback auto"),
    ("07", "Frontend React",             "Bilingue AR/FR, dark mode, visiteurs live"),
    ("08", "Tests & Resultats",          "Validation, metriques, KPIs"),
    ("09", "Demo & Liens",               "ChatBot + Dashboard Admin en ligne"),
    ("10", "Conclusion & Perspectives",  "Bilan et evolution future"),
]
col1, col2 = items[:5], items[5:]
anim_s2 = []
for i, (num, title, sub_) in enumerate(col1):
    y = Inches(1.28) + i * Inches(1.08)
    r_ = rect(s2, Inches(0.4), y + Inches(0.08), Inches(0.5), Inches(0.5), fill=G_MAIN)
    t_ = txt(s2, num, Inches(0.4), y + Inches(0.08), Inches(0.5), Inches(0.5),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    t2 = txt(s2, title, Inches(1.0), y + Inches(0.1), Inches(5.2), Inches(0.3),
             size=13, bold=True, color=G_MAIN)
    txt(s2, sub_, Inches(1.0), y + Inches(0.42), Inches(5.2), Inches(0.28),
        size=10, color=GRAY)
    anim_s2.append(r_.shape_id)
    anim_s2.append(t2.shape_id)
for i, (num, title, sub_) in enumerate(col2):
    y = Inches(1.28) + i * Inches(1.08)
    r_ = rect(s2, Inches(6.9), y + Inches(0.08), Inches(0.5), Inches(0.5), fill=G_MED)
    t2 = txt(s2, title, Inches(7.5), y + Inches(0.1), Inches(5.3), Inches(0.3),
             size=13, bold=True, color=G_MAIN)
    txt(s2, sub_, Inches(7.5), y + Inches(0.42), Inches(5.3), Inches(0.28),
        size=10, color=GRAY)
    anim_s2.append(r_.shape_id)
    anim_s2.append(t2.shape_id)
sep(s2, H - Inches(0.42))
add_auto_anim(s2, anim_s2, step_ms=120)
add_transition(s2, "fade")


# ======================================================
#  SLIDE 3 — CONTEXTE
# ======================================================
s3 = new_slide()
rect(s3, 0, 0, W, H, fill=BGLIGHT)
header(s3, "Contexte & Problematique", "Mauritanie | UNA | Besoin numerique")
footer(s3, 3)

stats = [("30 000+", "Etudiants", G_MAIN, WHITE),
         ("5",        "Facultes",  BLUE,   WHITE),
         ("12 000+",  "Questions/an", ORANGE, WHITE),
         ("1er",      "Chatbot univ. MR", GOLD, DARK)]
stat_ids = []
for i, (v, l, bg, fg) in enumerate(stats):
    x = Inches(0.4) + i * Inches(3.1)
    r_ = rect(s3, x, Inches(1.2), Inches(2.9), Inches(1.15), fill=bg)
    t_ = txt(s3, v, x, Inches(1.28), Inches(2.9), Inches(0.65),
             size=32, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txt(s3, l, x, Inches(1.93), Inches(2.9), Inches(0.35),
        size=11, color=fg, align=PP_ALIGN.CENTER)
    stat_ids.append(r_.shape_id)

rect(s3, Inches(0.4), Inches(2.55), Inches(5.9), Inches(4.2),
     fill=LRED, line=RED, lw=Pt(2))
rect(s3, Inches(0.4), Inches(2.55), Inches(5.9), Inches(0.48), fill=RED)
h_prob = txt(s3, "Problematique", Inches(0.6), Inches(2.59), Inches(5.5), Inches(0.4),
             size=15, bold=True, color=WHITE)
probs = ["Pas de reponse rapide hors heures de bureau",
         "Administration surchargee de questions repetitives",
         "Aucun chatbot universitaire en Mauritanie",
         "Barriere linguistique : besoin AR + FR"]
prob_ids = [h_prob.shape_id]
for i, p in enumerate(probs):
    t_ = txt(s3, "  " + p, Inches(0.6), Inches(3.12) + i * Inches(0.72),
             Inches(5.5), Inches(0.65), size=12, color=SLATE)
    prob_ids.append(t_.shape_id)

rect(s3, Inches(6.7), Inches(2.55), Inches(6.2), Inches(4.2),
     fill=G_PALE, line=G_MAIN, lw=Pt(2))
rect(s3, Inches(6.7), Inches(2.55), Inches(6.2), Inches(0.48), fill=G_MAIN)
h_sol = txt(s3, "Notre Solution", Inches(6.9), Inches(2.59), Inches(5.8), Inches(0.4),
            size=15, bold=True, color=WHITE)
sols = ["Chatbot IA bilingue disponible 24h/24, 7j/7",
        "Couvre 5 facultes + administration centrale",
        "300+ intentions avec detection automatique de langue",
        "Dashboard admin pour les questions sans reponse",
        "Deploye sur Railway - accessible partout"]
sol_ids = [h_sol.shape_id]
for i, s_ in enumerate(sols):
    t_ = txt(s3, "  " + s_, Inches(6.9), Inches(3.12) + i * Inches(0.72),
             Inches(5.8), Inches(0.65), size=12, color=SLATE)
    sol_ids.append(t_.shape_id)

add_auto_anim(s3, stat_ids + prob_ids + sol_ids, step_ms=180)
add_transition(s3, "fade")


# ======================================================
#  SLIDE 4 — ETAT DE L'ART
# ======================================================
s4 = new_slide()
rect(s4, 0, 0, W, H, fill=BGLIGHT)
header(s4, "Etat de l'Art", "Comparaison des solutions de chatbot existantes")
footer(s4, 4)

data4 = [
    ["Critere", "Dialogflow (Google)", "IBM Watson", "RASA Open Source", "Notre ChatBot UNA"],
    ["Open Source",        "Non",     "Non",     "OUI",    "OUI"],
    ["Langue Arabe",       "Limitee", "Limitee", "OUI",    "OUI (dialecte)"],
    ["Deploiement local",  "Non",     "Non",     "OUI",    "OUI (Railway)"],
    ["Cout",               "Payant",  "Payant",  "Gratuit","Gratuit"],
    ["Multi-tour",         "Oui",     "Oui",     "OUI",    "OUI (TEDPolicy)"],
    ["Connaissance metier","Externe", "Externe", "Interne","Interne (domain.yml)"],
    ["Dashboard admin",    "Basique", "Avance",  "Non",    "OUI (custom)"],
    ["Personnalisation",   "Limitee", "Limitee", "OUI",    "OUI (300+ intents)"],
]
cw4 = [Inches(2.4), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.4)]
tbl4 = s4.shapes.add_table(len(data4), 5,
       Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.44) * len(data4)).table
tbl4.first_row = True
for ci, cw in enumerate(cw4):
    tbl4.columns[ci].width = cw
for ri, row in enumerate(data4):
    for ci, cell_txt in enumerate(row):
        cell = tbl4.cell(ri, ci)
        cell.text = str(cell_txt)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = "Calibri"
        if ri == 0:
            run.font.size = Pt(11); run.font.bold = True
            run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = G_MAIN
        elif ci == 4:
            run.font.size = Pt(10); run.font.bold = True
            run.font.color.rgb = G_DARK
            cell.fill.solid(); cell.fill.fore_color.rgb = G_PALE
        else:
            run.font.size = Pt(10); run.font.color.rgb = SLATE
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else BGLIGHT
add_transition(s4, "fade")


# ======================================================
#  SLIDE 5 — ARCHITECTURE
# ======================================================
s5 = new_slide()
rect(s5, 0, 0, W, H, fill=BGLIGHT)
header(s5, "Architecture du Systeme", "3-tiers : React  Flask  Rasa  SQLite  Docker")
footer(s5, 5)

layers = [
    ("PRESENTATION", "React.js 18",
     ["Chat UI AR/FR", "Dark/Light Mode", "Visiteurs live", "Responsive"],
     BLUE, LBLUE, Inches(0.35)),
    ("METIER", "Flask API (Python)",
     ["10 endpoints REST", "Session management", "Fallback detection", "SQLite WAL"],
     G_MAIN, G_PALE, Inches(4.6)),
    ("DONNEES", "Rasa 3.x + SQLite",
     ["DIETClassifier", "TEDPolicy", "300+ intents", "Volume partage"],
     ORANGE, LORANGE, Inches(8.85)),
]
layer_ids = []
for (lbl, name, items_, fg, bg, x) in layers:
    r_ = rect(s5, x, Inches(1.2), Inches(3.9), Inches(5.5), fill=bg, line=fg, lw=Pt(2))
    rect(s5, x, Inches(1.2), Inches(3.9), Inches(0.42), fill=fg)
    t_lbl = txt(s5, lbl, x, Inches(1.22), Inches(3.9), Inches(0.38),
                size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s5, x + Inches(0.2), Inches(1.7), Inches(3.5), Inches(0.62), fill=fg)
    t_name = txt(s5, name, x + Inches(0.2), Inches(1.7), Inches(3.5), Inches(0.62),
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, it in enumerate(items_):
        txt(s5, "   " + it, x + Inches(0.3),
            Inches(2.45) + i * Inches(0.55), Inches(3.3), Inches(0.5),
            size=12, color=SLATE)
    layer_ids += [r_.shape_id, t_name.shape_id]

for xf in [Inches(4.25), Inches(8.5)]:
    txt(s5, "->", xf, Inches(3.75), Inches(0.38), Inches(0.4),
        size=20, bold=True, color=G_MAIN, align=PP_ALIGN.CENTER)
badge(s5, Inches(0.35), Inches(6.85), Inches(4.5), Inches(0.3),
      "Docker Compose - 3 services containerises", G_DARK, WHITE, size=10)
badge(s5, Inches(5.0), Inches(6.85), Inches(3.5), Inches(0.3),
      "SQLite WAL - volume partage", PURPLE, WHITE, size=10)

add_auto_anim(s5, layer_ids, step_ms=300)
add_transition(s5, "fade")


# ======================================================
#  SLIDE 6 — PIPELINE NLU
# ======================================================
s6 = new_slide()
rect(s6, 0, 0, W, H, fill=BGLIGHT)
header(s6, "Pipeline NLU - Rasa 3.x", "Traitement bilingue AR/FR")
footer(s6, 6)

comps = [
    ("Whitespace\nTokenizer",    BLUE,   LBLUE),
    ("Regex\nFeaturizer",        ORANGE, LORANGE),
    ("Count Vectors\nword+char", G_MAIN, G_PALE),
    ("DIET\nClassifier",         PURPLE, LPURPLE),
    ("Fallback\nClassifier",     RED,    LRED),
]
bw = Inches(2.35); bh = Inches(1.35); x0 = Inches(0.35)
comp_ids = []
for i, (name, fg, bg) in enumerate(comps):
    x = x0 + i * (bw + Inches(0.1))
    r_ = rect(s6, x, Inches(1.25), bw, bh, fill=bg, line=fg, lw=Pt(2))
    t_ = txt(s6, name, x, Inches(1.35), bw, Inches(0.7),
             size=13, bold=True, color=fg, align=PP_ALIGN.CENTER)
    comp_ids += [r_.shape_id, t_.shape_id]
    if i < len(comps) - 1:
        txt(s6, ">", x + bw + Inches(0.01), Inches(1.6), Inches(0.12), Inches(0.4),
            size=18, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

details = [
    ("Tokenisation", "Espaces, robuste arabe"),
    ("Patterns",     "Emails, nums, entites"),
    ("Vecteurs TF",  "n-grams 1-4 (char_wb)"),
    ("100 epochs",   "Intent + entites"),
    ("Seuil 0.3",    "Fallback automatique"),
]
for i, (_t, _b) in enumerate(details):
    x = x0 + i * (bw + Inches(0.1))
    txt(s6, _t, x + Inches(0.1), Inches(2.68), bw - Inches(0.2), Inches(0.3),
        size=10, bold=True, color=GRAY)
    txt(s6, _b, x + Inches(0.1), Inches(2.98), bw - Inches(0.2), Inches(0.45),
        size=9, color=SLATE)

txt(s6, "Politiques de dialogue", Inches(0.35), Inches(3.65), Inches(6), Inches(0.4),
    size=15, bold=True, color=G_MAIN)
sep(s6, Inches(4.08), G_LIGHT)
pols = [
    ("MemoizationPolicy", "Memorise chemins vus (max_history=5)", G_MAIN, G_PALE),
    ("TEDPolicy",         "Generalise dialogues inconnus (100 epochs)", BLUE, LBLUE),
    ("RulePolicy",        "Actions deterministiques (greetings...)", PURPLE, LPURPLE),
    ("UnexpecTEDIntent",  "Detecte intentions hors contexte", ORANGE, LORANGE),
]
pol_ids = []
for i, (pol, desc, fg, bg) in enumerate(pols):
    x = Inches(0.35) + i * Inches(3.22)
    r_ = rect(s6, x, Inches(4.15), Inches(3.1), Inches(1.1), fill=bg, line=fg, lw=Pt(1))
    txt(s6, pol, x + Inches(0.1), Inches(4.22), Inches(2.9), Inches(0.35),
        size=11, bold=True, color=fg)
    txt(s6, desc, x + Inches(0.1), Inches(4.58), Inches(2.9), Inches(0.55),
        size=9, color=SLATE)
    pol_ids.append(r_.shape_id)

r_lang = rect(s6, Inches(0.35), Inches(5.45), Inches(12.6), Inches(1.7),
              fill=G_PALE, line=G_MAIN, lw=Pt(1))
txt(s6, "Detection de langue : /[\\u0600-\\u06FF]/g  - si > 20% = Arabe (RTL)",
    Inches(0.55), Inches(5.6), Inches(11.5), Inches(0.4),
    size=13, bold=True, color=G_MAIN)
txt(s6, "char_wb n-gram (1 a 4) : robustesse aux fautes et au dialecte hassaniya",
    Inches(0.55), Inches(6.05), Inches(11.5), Inches(0.4),
    size=11, color=SLATE)

add_auto_anim(s6, comp_ids + pol_ids + [r_lang.shape_id], step_ms=180)
add_transition(s6, "fade")


# ======================================================
#  SLIDE 7 — SEQUENCE
# ======================================================
s7 = new_slide()
rect(s7, 0, 0, W, H, fill=BGLIGHT)
header(s7, "Diagramme de Sequence", "Flux complet d'un message utilisateur")
footer(s7, 7)

actors = [
    ("Utilisateur",  Inches(1.1), BLUE),
    ("React\nFrontend", Inches(3.4), G_MAIN),
    ("Flask\nBackend",  Inches(5.9), ORANGE),
    ("Rasa\nServer",    Inches(8.6), PURPLE),
    ("SQLite",          Inches(11.3), RED),
]
actor_ids = []
for name, x, color in actors:
    r_ = rect(s7, x - Inches(0.65), Inches(1.2), Inches(1.3), Inches(0.6), fill=color)
    txt(s7, name, x - Inches(0.65), Inches(1.2), Inches(1.3), Inches(0.6),
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s7, x - Pt(1), Inches(1.8), Pt(2), Inches(4.9), fill=LGRAY)
    actor_ids.append(r_.shape_id)

msgs = [
    (Inches(1.1),  Inches(3.4),  "1. Envoie message",   BLUE),
    (Inches(3.4),  Inches(5.9),  "2. POST /chat",        G_MAIN),
    (Inches(5.9),  Inches(8.6),  "3. Webhook Rasa",      ORANGE),
    (Inches(8.6),  Inches(5.9),  "4. JSON reponse",      ORANGE),
    (Inches(5.9),  Inches(11.3), "5. Fallback -> Save",  RED),
    (Inches(5.9),  Inches(3.4),  "6. Reponse JSON",      G_MAIN),
    (Inches(3.4),  Inches(1.1),  "7. Affiche reponse",   BLUE),
]
msg_ids = []
for i, (x1, x2, msg, c) in enumerate(msgs):
    yp = Inches(2.0) + i * Inches(0.72)
    xs = min(x1, x2) + Inches(0.04)
    xe = max(x1, x2) - Inches(0.04)
    r_ = rect(s7, xs, yp, xe - xs, Pt(2), fill=c)
    arrow = ">" if x2 > x1 else "<"
    ax = xe - Inches(0.18) if x2 > x1 else xs - Inches(0.04)
    txt(s7, arrow, ax, yp - Inches(0.18), Inches(0.2), Inches(0.25),
        size=12, bold=True, color=c)
    mx = (xs + xe) / 2 - Inches(0.9)
    txt(s7, msg, mx, yp - Inches(0.25), Inches(2.0), Inches(0.22),
        size=8, color=c, align=PP_ALIGN.CENTER)
    msg_ids.append(r_.shape_id)

rect(s7, Inches(0.35), Inches(7.1), W - Inches(0.7), Inches(0.26),
     fill=G_PALE, line=G_LIGHT, lw=Pt(1))
txt(s7, "En cas de fallback : question auto-sauvegardee en SQLite - aucune perte de donnee",
    Inches(0.55), Inches(7.13), W - Inches(1.0), Inches(0.2),
    size=9, color=G_MAIN, italic=True, align=PP_ALIGN.CENTER)

add_auto_anim(s7, actor_ids + msg_ids, step_ms=150)
add_transition(s7, "fade")


# ======================================================
#  SLIDE 8 — FLASK BACKEND
# ======================================================
s8 = new_slide()
rect(s8, 0, 0, W, H, fill=BGLIGHT)
header(s8, "Backend Flask - API REST", "Python 3  SQLite WAL  CORS  10 endpoints")
footer(s8, 8)

ep = [
    ["Methode", "Route", "Description"],
    ["POST",   "/chat",                         "Envoie message a Rasa, retourne la reponse"],
    ["GET",    "/health",                        "Status API (Flask + Rasa)"],
    ["GET",    "/ping",                          "Keepalive ping"],
    ["GET",    "/active-users",                  "Sessions actives en temps reel"],
    ["GET",    "/unknown-questions",             "Liste questions sans reponse"],
    ["POST",   "/unknown-questions/<id>/reply",  "Ajouter reponse admin"],
    ["DELETE", "/unknown-questions/<id>",        "Supprimer une question"],
    ["POST",   "/save-unknown-question",         "Sauvegarder (depuis Rasa actions)"],
    ["POST",   "/register-session",              "Enregistrer session visiteur"],
    ["POST",   "/deregister-session",            "Desenregistrer session"],
]
cw8 = [Inches(1.3), Inches(3.9), Inches(7.1)]
tbl8 = s8.shapes.add_table(len(ep), 3,
       Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.42) * len(ep)).table
for ci, cw in enumerate(cw8):
    tbl8.columns[ci].width = cw
mc = {"POST": G_MAIN, "GET": BLUE, "DELETE": RED, "PUT": ORANGE}
for ri, row in enumerate(ep):
    for ci, ct in enumerate(row):
        cell = tbl8.cell(ri, ci)
        cell.text = str(ct)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci < 2 else PP_ALIGN.LEFT
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = "Calibri"
        if ri == 0:
            run.font.size = Pt(11); run.font.bold = True
            run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = G_MAIN
        elif ci == 0:
            run.font.size = Pt(10); run.font.bold = True
            run.font.color.rgb = mc.get(row[0], GRAY)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
        else:
            run.font.size = Pt(10); run.font.color.rgb = SLATE
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else BGLIGHT

badge(s8, Inches(0.4), Inches(6.92), Inches(6.0), Inches(0.3),
      "SQLite WAL - lectures et ecritures simultanees Flask + Rasa", G_DARK, WHITE, size=10)
add_transition(s8, "fade")


# ======================================================
#  SLIDE 9 — REACT FRONTEND
# ======================================================
s9 = new_slide()
rect(s9, 0, 0, W, H, fill=BGLIGHT)
header(s9, "Frontend React.js 18", "Interface bilingue AR/FR - Dark/Light Mode - Visiteurs live")
footer(s9, 9)

feats = [
    ("Detection langue",  "Regex Unicode /[\\u0600-\\u06FF]/g\n> 20% = Arabe, RTL auto",  BLUE,   LBLUE),
    ("Mode sombre",       "Dark/Light theme\nDark par defaut\nBouton toggle",              G_MAIN, G_PALE),
    ("Visiteurs live",    "Ping Flask /30s\nCompteur affiche\nFormatage 1K / 10K",         ORANGE, LORANGE),
    ("Sante systeme",     "Health check /20s\nIndicateur vert/rouge\nReconnexion auto",    RED,    LRED),
    ("Interface",         "Tajawal (police AR)\nLucide React icons\nBulles de chat",       PURPLE, LPURPLE),
    ("Dashboard",         "Route /dashboard\nQuestions inconnues\nReplies admin",          GOLD,   LGOLD),
]
feat_ids = []
for i, (t, b, fg, bg) in enumerate(feats):
    col = i % 3; row = i // 3
    x = Inches(0.4) + col * Inches(4.25)
    y = Inches(1.3) + row * Inches(2.05)
    card(s9, x, y, Inches(4.0), Inches(1.9))
    r_ = rect(s9, x, y, Inches(4.0), Inches(0.45), fill=fg)
    txt(s9, t, x + Inches(0.15), y + Inches(0.07), Inches(3.7), Inches(0.32),
        size=13, bold=True, color=WHITE)
    txt(s9, b, x + Inches(0.15), y + Inches(0.55), Inches(3.7), Inches(1.25),
        size=11, color=SLATE)
    feat_ids.append(r_.shape_id)

add_auto_anim(s9, feat_ids, step_ms=250)
add_transition(s9, "fade")


# ======================================================
#  SLIDE 10 — DOCKER & DEPLOY
# ======================================================
s10 = new_slide()
rect(s10, 0, 0, W, H, fill=BGLIGHT)
header(s10, "Docker & Deploiement", "3 services containerises - Railway Cloud")
footer(s10, 10)

containers = [
    ("react-frontend", ":80",
     "Nginx\nSert l'app React\nBuild-time env vars\nHTTPS auto",
     BLUE, LBLUE),
    ("flask-backend", ":5000",
     "Python 3.11\nAPI REST\nVolume /data SQLite\nCORS actif",
     G_MAIN, G_PALE),
    ("rasa-server", ":5005\n:5055",
     "Rasa 3.x\nModele NLU\nActions server\nVolume modele",
     ORANGE, LORANGE),
]
cont_ids = []
for i, (name, port, desc, fg, bg) in enumerate(containers):
    x = Inches(0.45) + i * Inches(4.25)
    r_ = rect(s10, x, Inches(1.25), Inches(4.0), Inches(3.0), fill=bg, line=fg, lw=Pt(2))
    rect(s10, x, Inches(1.25), Inches(4.0), Inches(0.55), fill=fg)
    txt(s10, name, x + Inches(0.12), Inches(1.29), Inches(2.6), Inches(0.45),
        size=14, bold=True, color=WHITE)
    badge(s10, x + Inches(2.75), Inches(1.33), Inches(1.1), Inches(0.38),
          port, WHITE, fg, size=10, bold=True)
    txt(s10, desc, x + Inches(0.18), Inches(1.88), Inches(3.65), Inches(1.85),
        size=11, color=SLATE)
    cont_ids.append(r_.shape_id)

conns = [
    ("React -> Flask",  "/chat, /health, /ping, /active-users",    BLUE),
    ("Flask -> Rasa",   "webhooks/rest/webhook (POST)",             G_MAIN),
    ("Rasa -> SQLite",  "Volume /data partage - ecriture WAL",      ORANGE),
    ("Flask -> SQLite", "Lecture/ecriture questions inconnues",      RED),
]
for i, (t, d, c) in enumerate(conns):
    y = Inches(4.5) + i * Inches(0.5)
    rect(s10, Inches(0.4), y + Inches(0.13), Inches(0.3), Inches(0.22), fill=c)
    txt(s10, t + " :", Inches(0.8), y, Inches(2.8), Inches(0.42),
        size=12, bold=True, color=c)
    txt(s10, d, Inches(3.65), y, Inches(6.5), Inches(0.42), size=11, color=SLATE)

rect(s10, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.55),
     fill=G_PALE, line=G_MAIN, lw=Pt(1))
txt(s10, "Deploye sur Railway.app  |  HTTPS auto  |  Volume /data persistant  |  Env vars injectees",
    Inches(0.6), Inches(6.73), Inches(9.5), Inches(0.35), size=11, color=SLATE)
btn(s10, " Voir le deploiement ", URL_CHATBOT,
    Inches(10.2), Inches(6.65), Inches(2.7), Inches(0.52),
    bg=G_MAIN, fg=WHITE, size=11)

add_auto_anim(s10, cont_ids, step_ms=350)
add_transition(s10, "fade")


# ======================================================
#  SLIDE 11 — TESTS
# ======================================================
s11 = new_slide()
rect(s11, 0, 0, W, H, fill=BGLIGHT)
header(s11, "Tests & Validation", "Tests unitaires, d'integration et de charge")
footer(s11, 11)

data11 = [
    ["Type", "Composant", "Scenario", "Resultat"],
    ["Unite",       "Flask /chat",    "Message valide -> reponse Rasa",     "PASS"],
    ["Unite",       "Flask /health",  "Status ok quand Rasa actif",         "PASS"],
    ["Unite",       "Rasa NLU",       "Intent reconnu seuil > 0.7",         "PASS"],
    ["Unite",       "Fallback",       "Intent inconnu -> SQLite save",       "PASS"],
    ["Integration", "React->Flask",  "Envoi message + affichage reponse",  "PASS"],
    ["Integration", "Flask->Rasa",   "Webhook + retour JSON",              "PASS"],
    ["Integration", "Rasa->SQLite",  "Question inconnue -> DB partagee",   "PASS"],
    ["Charge",      "10 users simult","Temps reponse < 2s",                "PASS"],
    ["Bilingue",    "AR detection",  "> 20% arabes -> RTL correct",        "PASS"],
    ["Admin",       "/dashboard",    "Lecture + reponse questions",        "PASS"],
]
cw11 = [Inches(1.7), Inches(2.2), Inches(6.4), Inches(1.8)]
tbl11 = s11.shapes.add_table(len(data11), 4,
        Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.44) * len(data11)).table
for ci, cw in enumerate(cw11):
    tbl11.columns[ci].width = cw
for ri, row in enumerate(data11):
    for ci, ct in enumerate(row):
        cell = tbl11.cell(ri, ci)
        cell.text = str(ct)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = "Calibri"
        if ri == 0:
            run.font.size = Pt(11); run.font.bold = True
            run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = G_MAIN
        elif ci == 3:
            run.font.size = Pt(10); run.font.bold = True
            run.font.color.rgb = G_DARK
            cell.fill.solid(); cell.fill.fore_color.rgb = G_PALE
        else:
            run.font.size = Pt(10); run.font.color.rgb = SLATE
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else BGLIGHT

badge(s11, Inches(7.5), Inches(6.9), Inches(5.3), Inches(0.32),
      "Taux de reussite : 100%  (10 / 10 scenarios PASS)", G_MAIN, WHITE, size=11)
add_transition(s11, "fade")


# ======================================================
#  SLIDE 12 — DASHBOARD ADMIN
# ======================================================
s12 = new_slide()
rect(s12, 0, 0, W, H, fill=BGLIGHT)
header(s12, "Panneau Administration", "Gestion des questions inconnues + reponses admin")
footer(s12, 12)

txt(s12, "Acces direct :", Inches(0.4), Inches(1.22), Inches(2.2), Inches(0.42),
    size=13, bold=True, color=G_MAIN)
btn(s12, "  Ouvrir le Dashboard Admin  ", URL_ADMIN,
    Inches(2.7), Inches(1.18), Inches(5.0), Inches(0.52),
    bg=G_LIGHT, fg=G_DARK, size=13)
txt_link(s12, URL_ADMIN, URL_ADMIN,
         Inches(7.9), Inches(1.28), Inches(5.2), Inches(0.32),
         size=10, color=BLUE)

feats12 = [
    ("1. Consultation",
     "Liste toutes les questions non repondues, triees par date (plus recentes en premier).",
     G_MAIN, G_PALE),
    ("2. Repondre",
     "L'admin saisit une reponse directement. Stockee en SQLite avec horodatage replied_at.",
     BLUE, LBLUE),
    ("3. Supprimer",
     "Supprimer une question resolue en un clic - DELETE /unknown-questions/<id>.",
     RED, LRED),
    ("4. Statistiques",
     "Compteur total, questions sans reponse, taux de resolution, export possible.",
     ORANGE, LORANGE),
]
feat12_ids = []
for i, (t, b, fg, bg) in enumerate(feats12):
    col = i % 2; row = i // 2
    x = Inches(0.4) + col * Inches(6.4)
    y = Inches(1.9) + row * Inches(2.15)
    card(s12, x, y, Inches(6.1), Inches(2.0))
    r_ = rect(s12, x, y, Inches(6.1), Inches(0.5), fill=fg)
    txt(s12, t, x + Inches(0.18), y + Inches(0.08), Inches(5.7), Inches(0.35),
        size=14, bold=True, color=WHITE)
    txt(s12, b, x + Inches(0.18), y + Inches(0.62), Inches(5.7), Inches(1.2),
        size=12, color=SLATE)
    feat12_ids.append(r_.shape_id)

rect(s12, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.85),
     fill=G_PALE, line=G_MAIN, lw=Pt(1))
txt(s12, "Route React : /dashboard  |  API : GET /unknown-questions + POST /reply + DELETE",
    Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.45),
    size=11, color=SLATE, align=PP_ALIGN.CENTER)

add_auto_anim(s12, feat12_ids, step_ms=300)
add_transition(s12, "fade")


# ======================================================
#  SLIDE 13 — DEMO LIVE
# ======================================================
s13 = new_slide()
rect(s13, 0, 0, W, H, fill=G_DARK)
for i in range(20):
    rect(s13, Inches(i * 0.68), 0, Inches(0.3), H,
         fill=RGBColor(0x12, 0x4a, 0x2c))

t_demo = txt(s13, "Demo Live", Inches(0.5), Inches(0.35), Inches(9), Inches(0.65),
             size=34, bold=True, color=WHITE)
rect(s13, Inches(0.5), Inches(1.0), Inches(3.5), Pt(3), fill=G_LIGHT)

# Chatbot card
r_chat = rect(s13, Inches(0.45), Inches(1.12), Inches(5.9), Inches(5.7),
              fill=RGBColor(0x12, 0x35, 0x24), line=G_LIGHT, lw=Pt(2))
rect(s13, Inches(0.45), Inches(1.12), Inches(5.9), Inches(0.62), fill=G_MAIN)
txt(s13, "ChatBot UNA - Interface utilisateur",
    Inches(0.62), Inches(1.17), Inches(5.55), Inches(0.52),
    size=13, bold=True, color=WHITE)

demo_msgs = [
    (" Quels sont les filieres de la FST ?",              G_MED, WHITE, Inches(4.6)),
    (" La FST propose : Informatique, Maths, Physique...", BLUE,  WHITE, Inches(0.65)),
    (" ma hi takhassousat kulliyyat al-iqtisad ?",         G_MED, WHITE, Inches(4.6)),
    (" tuqaddimu kulliyyat al-iqtisad takhassousat...",    BLUE,  WHITE, Inches(0.65)),
    (" Comment s'inscrire a la FST ?",                     G_MED, WHITE, Inches(4.6)),
    (" Rendez-vous sur fst-preinscription.netlify.app",    BLUE,  WHITE, Inches(0.65)),
]
msg_ids_13 = []
for i, (msg, bg, fg, xp) in enumerate(demo_msgs):
    b_ = badge(s13, xp, Inches(1.88) + i * Inches(0.68), Inches(4.8), Inches(0.55),
               msg, bg, fg, size=10, bold=False)
    msg_ids_13.append(b_.shape_id)

b_open_chat = btn(s13, "  Ouvrir le ChatBot  ", URL_CHATBOT,
                  Inches(0.62), Inches(6.4), Inches(5.55), Inches(0.65),
                  bg=G_LIGHT, fg=G_DARK, size=14)

# Admin card
r_admin = rect(s13, Inches(7.0), Inches(1.12), Inches(5.9), Inches(5.7),
               fill=RGBColor(0x12, 0x35, 0x24), line=G_LIGHT, lw=Pt(2))
rect(s13, Inches(7.0), Inches(1.12), Inches(5.9), Inches(0.62), fill=G_MED)
txt(s13, "Dashboard Admin - Questions sans reponse",
    Inches(7.18), Inches(1.17), Inches(5.55), Inches(0.52),
    size=13, bold=True, color=WHITE)

admin_rows = [
    ("Quand commence la session 2025 ?", "2025-05-01"),
    ("Comment imprimer mon releve ?",    "2025-05-02"),
    ("Lien du site officiel FST ?",      "2025-05-03"),
    ("Y a-t-il des bourses disponibles?","2025-05-04"),
]
for i, (q, dt) in enumerate(admin_rows):
    yp = Inches(1.88) + i * Inches(0.82)
    rect(s13, Inches(7.1), yp, Inches(5.7), Inches(0.7),
         fill=RGBColor(0x1a, 0x3d, 0x2b), line=G_MED, lw=Pt(1))
    txt(s13, q, Inches(7.22), yp + Inches(0.05), Inches(4.5), Inches(0.32),
        size=10, color=WHITE)
    txt(s13, dt, Inches(7.22), yp + Inches(0.38), Inches(2.5), Inches(0.24),
        size=9, color=G_TEXT)

b_open_admin = btn(s13, "  Ouvrir le Dashboard  ", URL_ADMIN,
                   Inches(7.18), Inches(6.4), Inches(5.55), Inches(0.65),
                   bg=G_MED, fg=WHITE, size=14)

add_auto_anim(s13,
              [t_demo.shape_id, r_chat.shape_id, r_admin.shape_id] +
              msg_ids_13 +
              [b_open_chat.shape_id, b_open_admin.shape_id],
              step_ms=200)
add_transition(s13, "fade")


# ======================================================
#  SLIDE 14 — RESULTATS KPIs
# ======================================================
s14 = new_slide()
rect(s14, 0, 0, W, H, fill=BGLIGHT)
header(s14, "Resultats & Metriques", "Performances mesures sur le systeme deploye")
footer(s14, 14)

kpis = [
    ("300+",  "Intentions NLU",  G_MAIN, WHITE),
    ("98%",   "Precision NLU",   BLUE,   WHITE),
    ("<1.5s", "Temps reponse",   ORANGE, WHITE),
    ("24/7",  "Disponibilite",   G_DARK, WHITE),
    ("2",     "Langues AR+FR",   PURPLE, WHITE),
    ("5",     "Facultes",        RED,    WHITE),
]
kw = Inches(2.0)
kpi_ids = []
for i, (v, l, bg, fg) in enumerate(kpis):
    x = Inches(0.4) + i * Inches(2.15)
    r_ = rect(s14, x, Inches(1.25), kw, Inches(1.62), fill=bg)
    txt(s14, v, x, Inches(1.32), kw, Inches(0.85),
        size=34, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txt(s14, l, x, Inches(2.18), kw, Inches(0.55),
        size=11, color=fg, align=PP_ALIGN.CENTER)
    kpi_ids.append(r_.shape_id)

txt(s14, "Contributions principales", Inches(0.4), Inches(3.12), Inches(8), Inches(0.45),
    size=17, bold=True, color=G_MAIN)
sep(s14, Inches(3.58), G_LIGHT)

contribs = [
    ("1er chatbot universitaire MR",
     "Concu et deploye entierement par des etudiants mauritaniens"),
    ("Stack full-stack professionnel",
     "React 18 + Flask + Rasa 3.x + Docker + Railway : niveau production"),
    ("Bilinguisme natif AR/FR",
     "Detection automatique, RTL/LTR dynamique, dialecte hassaniya"),
    ("Persistance multi-niveaux",
     "3 mecanismes : SQLite -> Flask API -> JSON - aucune perte possible"),
]
contrib_ids = []
for i, (t, b) in enumerate(contribs):
    col = i % 2; row = i // 2
    x = Inches(0.4) + col * Inches(6.45)
    y = Inches(3.72) + row * Inches(1.52)
    r_ = rect(s14, x, y, Inches(0.08), Inches(1.3), fill=G_LIGHT)
    txt(s14, t, x + Inches(0.22), y + Inches(0.05), Inches(5.9), Inches(0.38),
        size=13, bold=True, color=G_MAIN)
    txt(s14, b, x + Inches(0.22), y + Inches(0.48), Inches(5.9), Inches(0.75),
        size=11, color=SLATE)
    contrib_ids.append(r_.shape_id)

add_auto_anim(s14, kpi_ids + contrib_ids, step_ms=200)
add_transition(s14, "fade")


# ======================================================
#  SLIDE 15 — BASE DE CONNAISSANCES
# ======================================================
s15 = new_slide()
rect(s15, 0, 0, W, H, fill=BGLIGHT)
header(s15, "Base de Connaissances", "300+ intentions couvrant 5 facultes + administration")
footer(s15, 15)

data15 = [
    ["Faculte / Domaine",         "Code",  "Intents", "Themes principaux"],
    ["Sciences & Techniques",     "FST",   "85+",
     "Filieres, inscription, resultats, examens, contacts"],
    ["Economie & Gestion",        "FEG",   "55+",
     "Matieres, stages, debouches, modalites"],
    ["Droit & Sciences Po",       "FSJP",  "22+",
     "Droits, licences, inscriptions, cursus"],
    ["Lettres & Sc. Humaines",    "FLSH",  "17+",
     "Langues, histoire, geographie, arts"],
    ["Medecine & Pharmacie",      "FMPOS", "8+",
     "Medecine, pharmacie, chirurgie dentaire"],
    ["Administration Centrale",   "ADM",   "55+",
     "Scolarite, calendrier, logement, bourse, contacts"],
    ["TOTAL (toutes facultes)",   "-",     "300+",
     "Couverture complete de l'UNA - AR + FR"],
]
cw15 = [Inches(3.5), Inches(1.4), Inches(1.7), Inches(5.7)]
make_table(s15, Inches(0.4), Inches(1.25), Inches(12.5), data15, cw15)
txt(s15, "Chaque intention possede des variantes en arabe ET en francais avec synonymes et entites",
    Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.45),
    size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
add_transition(s15, "fade")


# ======================================================
#  SLIDE 16 — CONCLUSION
# ======================================================
s16 = new_slide()
rect(s16, 0, 0, W, H, fill=BGLIGHT)
header(s16, "Conclusion", "Bilan du Projet de Fin d'Etudes")
footer(s16, 16)

r_obj = rect(s16, Inches(0.4), Inches(1.25), Inches(5.9), Inches(5.5),
             fill=G_PALE, line=G_MAIN, lw=Pt(2))
rect(s16, Inches(0.4), Inches(1.25), Inches(5.9), Inches(0.5), fill=G_MAIN)
txt(s16, "Objectifs atteints", Inches(0.6), Inches(1.3), Inches(5.5), Inches(0.4),
    size=15, bold=True, color=WHITE)
objs = ["ChatBot bilingue AR/FR deploye et fonctionnel",
        "Couverture des 5 facultes de l'UNA",
        "300+ intentions avec detection de fallback",
        "Dashboard admin operationnel",
        "Persistance SQLite sur 3 niveaux",
        "Deploiement cloud Docker + Railway",
        "Interface dark/light mode responsive",
        "1er chatbot universitaire en Mauritanie"]
for i, o in enumerate(objs):
    txt(s16, "   " + o, Inches(0.6), Inches(1.88) + i * Inches(0.54),
        Inches(5.5), Inches(0.5), size=12, color=SLATE)

r_persp = rect(s16, Inches(6.75), Inches(1.25), Inches(6.1), Inches(5.5),
               fill=G_PALE, line=G_MED, lw=Pt(2))
rect(s16, Inches(6.75), Inches(1.25), Inches(6.1), Inches(0.5), fill=G_MED)
txt(s16, "Perspectives", Inches(6.95), Inches(1.3), Inches(5.7), Inches(0.4),
    size=15, bold=True, color=WHITE)
persp = [
    ("Integration LLM",      "GPT-4/Claude pour questions complexes"),
    ("Version mobile",       "Application Android/iOS"),
    ("Reconnaissance vocale","Saisie et reponse audio AR/FR"),
    ("Apprentissage auto",   "Reentranement depuis le dashboard"),
    ("Multi-universite",     "Extension a d'autres universites"),
    ("API publique",         "Ouverture pour applications tierces"),
    ("Analyse semantique",   "Recherche floue et suggestions"),
]
for i, (t, b) in enumerate(persp):
    txt(s16, "   " + t + " :", Inches(6.95), Inches(1.88) + i * Inches(0.54),
        Inches(2.4), Inches(0.5), size=12, bold=True, color=G_MED)
    txt(s16, b, Inches(9.35), Inches(1.88) + i * Inches(0.54),
        Inches(3.3), Inches(0.5), size=12, color=SLATE)

add_auto_anim(s16, [r_obj.shape_id, r_persp.shape_id], step_ms=400)
add_transition(s16, "fade")


# ======================================================
#  SLIDE 17 — FACULTES
# ======================================================
s17 = new_slide()
rect(s17, 0, 0, W, H, fill=BGLIGHT)
header(s17, "Universite de Nouakchott", "5 Facultes - 30 000+ etudiants")
footer(s17, 17)

facs = [
    ("FST",   "Sciences &\nTechniques",    "Informatique, Maths,\nPhysique, Chimie", "85+ intents",  G_MAIN, G_PALE),
    ("FEG",   "Economie &\nGestion",       "Commerce, Finance,\nMarketing, RH",      "55+ intents",  BLUE,   LBLUE),
    ("FSJP",  "Droit &\nSc. Politiques",   "Droit prive, public,\nSciences Po",      "22+ intents",  PURPLE, LPURPLE),
    ("FLSH",  "Lettres &\nSc. Humaines",   "Arabe, Francais,\nHistoire, Geo",        "17+ intents",  ORANGE, LORANGE),
    ("FMPOS", "Medecine &\nPharmacie",     "Medecine, Pharmacie,\nChirurgie dent.",  "8+ intents",   RED,    LRED),
]
fw = Inches(2.35)
fac_ids = []
for i, (code, name, desc, intents, fg, bg) in enumerate(facs):
    x = Inches(0.3) + i * (fw + Inches(0.2))
    r_ = rect(s17, x, Inches(1.22), fw, Inches(0.62), fill=fg)
    txt(s17, code, x, Inches(1.22), fw, Inches(0.62),
        size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card(s17, x, Inches(1.84), fw, Inches(4.0), fill=bg, border=fg)
    txt(s17, name, x + Inches(0.1), Inches(1.96), fw - Inches(0.2), Inches(0.7),
        size=13, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txt(s17, desc, x + Inches(0.1), Inches(2.72), fw - Inches(0.2), Inches(1.1),
        size=10, color=SLATE, align=PP_ALIGN.CENTER)
    badge(s17, x + Inches(0.35), Inches(3.95), fw - Inches(0.7), Inches(0.32),
          intents, fg, WHITE, size=10)
    fac_ids.append(r_.shape_id)

txt(s17, "Administration Centrale : 55+ intents - Scolarite, Calendrier, Logement, Bourse",
    Inches(0.4), Inches(6.1), W - Inches(0.8), Inches(0.45),
    size=12, bold=True, color=G_MAIN, align=PP_ALIGN.CENTER)
badge(s17, Inches(4.2), Inches(6.65), Inches(4.8), Inches(0.45),
      "Total : 300+ intents  |  Bilingue AR + FR", G_DARK, WHITE, size=12)

add_auto_anim(s17, fac_ids, step_ms=280)
add_transition(s17, "fade")


# ======================================================
#  SLIDE 18 — MERCI
# ======================================================
s18 = new_slide()
rect(s18, 0, 0, W, H, fill=G_DARK)
for i in range(12):
    rect(s18, Inches(8.8 + i * 0.42), 0, Inches(0.28), H,
         fill=RGBColor(0x1a, 0x5c, 0x35))
rect(s18, 0, Inches(2.6), W, Inches(0.05), fill=G_LIGHT)
rect(s18, 0, Inches(5.45), W, Inches(0.05), fill=G_LIGHT)

t_merci  = txt(s18, "Merci pour votre attention !",
               Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.95),
               size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
t_una    = txt(s18, "ChatBot UNA  -  Universite de Nouakchott",
               Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.5),
               size=18, color=G_TEXT, italic=True, align=PP_ALIGN.CENTER)
t_liens  = txt(s18, "Liens du projet :",
               Inches(1.5), Inches(2.85), Inches(10.3), Inches(0.5),
               size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
b_chat   = btn(s18, "  ChatBot en ligne  ", URL_CHATBOT,
               Inches(1.5), Inches(3.55), Inches(4.5), Inches(0.78),
               bg=G_LIGHT, fg=G_DARK, size=15)
b_adm    = btn(s18, "  Dashboard Admin  ", URL_ADMIN,
               Inches(7.3), Inches(3.55), Inches(4.5), Inches(0.78),
               bg=G_MED, fg=WHITE, size=15)

tb_url = s18.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.35))
tf_url = tb_url.text_frame
p_url  = tf_url.paragraphs[0]
p_url.alignment = PP_ALIGN.CENTER
r_url  = p_url.add_run()
r_url.text = URL_CHATBOT
r_url.font.size = Pt(12)
r_url.font.color.rgb = G_TEXT
r_url.font.name = "Calibri"
try:
    r_url.hyperlink.address = URL_CHATBOT
except Exception:
    pass

t_qst = txt(s18, "Questions & Discussion",
            Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.62),
            size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
t_stack = txt(s18, "Rasa 3.x  |  Flask  |  React.js 18  |  SQLite WAL  |  Docker  |  Railway",
              Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
              size=13, color=G_TEXT, italic=True, align=PP_ALIGN.CENTER)

add_auto_anim(s18,
              [t_merci.shape_id, t_una.shape_id, t_liens.shape_id,
               b_chat.shape_id, b_adm.shape_id, t_qst.shape_id, t_stack.shape_id],
              step_ms=220)
add_transition(s18, "fade")


# ======================================================
#  SAUVEGARDE
# ======================================================
OUT = r"c:\Users\hp\Desktop\Presentation_PFE_ChatBot_UNA.pptx"
prs.save(OUT)
print("Fichier genere : " + OUT)
print("Slides : {}".format(len(prs.slides)))
print("ChatBot : " + URL_CHATBOT)
print("Admin   : " + URL_ADMIN)
